"""PowerPoint (.pptx) extractor."""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation

from extractors.base import Extractor

logger = logging.getLogger(__name__)


class PptxExtractor(Extractor):
    """Extract raw text from .pptx files.

    Iterates over slides → shapes → text frames → paragraphs.
    """

    def extract(self, file_path: str | Path) -> str:
        prs = Presentation(str(file_path))
        parts: list[str] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)

                # Also handle tables inside slides
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [
                            cell.text.strip()
                            for cell in row.cells
                            if cell.text.strip()
                        ]
                        if cells:
                            slide_texts.append("\t".join(cells))

            if slide_texts:
                parts.append(f"[Slide {slide_idx}]")
                parts.extend(slide_texts)

        return "\n".join(parts)
